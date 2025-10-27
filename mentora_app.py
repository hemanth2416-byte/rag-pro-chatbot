# 🤖 Mentora v3.4.1 — AI Teaching Assistant with Voice, ChromaDB (HTTP Mode) + FAISS Fallback
# Author: Hemanth Kumar

import streamlit as st
import fitz, requests, json, numpy as np, os, re, tempfile, faiss
import docx, pandas as pd
from bs4 import BeautifulSoup
from pptx import Presentation
from PIL import Image
import pytesseract
import speech_recognition as sr
from pydub import AudioSegment
from moviepy.editor import VideoFileClip
from fpdf import FPDF
import plotly.express as px
from collections import Counter
from gtts import gTTS
import whisper
from chromadb import HttpClient

# ---------- Config ----------
OLLAMA_URL = os.getenv("OLLAMA_URL", "http://localhost:11434")
CHROMA_HOST, CHROMA_PORT = "localhost", 8001
SESSIONS_PATH, EXPORTS_PATH = "data/sessions", "exports"
os.makedirs(EXPORTS_PATH, exist_ok=True)
os.makedirs(SESSIONS_PATH, exist_ok=True)

st.set_page_config(page_title="🤖 Mentora | AI Teaching Assistant", layout="wide")

# ---------- Theme ----------
theme = st.sidebar.radio("🎨 Theme", ["🌞 Day Mode", "🌙 Night Mode"])
if theme == "🌞 Day Mode":
    bg = "linear-gradient(135deg,#f9fbff 0%,#d6e6ff 100%)"
    tcol = "#0a1a2f"; banner = "linear-gradient(90deg,#1748b1,#2e8bff)"
else:
    bg = "linear-gradient(135deg,#0f2027,#203a43,#2c5364)"
    tcol = "#ffffff"; banner = "linear-gradient(90deg,#0a84ff,#00274d)"

st.markdown(f"""
<link href="https://fonts.googleapis.com/css2?family=Poppins:wght@400;600&display=swap" rel="stylesheet">
<style>
body,html,.main{{font-family:'Poppins',sans-serif;background:{bg};color:{tcol};}}
section[data-testid="stSidebar"]{{background:#1748b1;color:white !important;}}
.title-banner{{background:{banner};border-radius:12px;padding:15px;text-align:center;color:white;margin-bottom:20px;}}
.footer{{text-align:center;color:gray;font-size:14px;margin-top:30px;}}
</style>
""", unsafe_allow_html=True)

st.markdown("""
<div class='title-banner'>
  <h2>🤖 Mentora — Your AI Teaching Assistant</h2>
  <p>Upload, speak, and learn with intelligent voice & text interaction.</p>
</div>
""", unsafe_allow_html=True)

# ---------- Persona ----------
persona = st.sidebar.selectbox(
    "🎙️ Teaching Style",
    ["Friendly Mentor", "Strict Professor", "Motivational Coach", "Coding Tutor"],
    index=0
)
PERSONA_STYLE = {
    "Friendly Mentor": "Be warm, encouraging, and explain step by step.",
    "Strict Professor": "Be concise and emphasize clarity and definitions.",
    "Motivational Coach": "Be positive, energetic, and supportive.",
    "Coding Tutor": "Focus on clear code examples and real-world practice."
}

# ---------- Whisper ----------
@st.cache_resource
def load_whisper_model():
    return whisper.load_model("small")

# ---------- Extraction ----------
def extract_text_from_image(f): return pytesseract.image_to_string(Image.open(f))
def extract_text_from_audio(f):
    sound = AudioSegment.from_file(f); sound.export("temp.wav", format="wav")
    r = sr.Recognizer()
    with sr.AudioFile("temp.wav") as s: a = r.record(s)
    try: return r.recognize_google(a)
    except Exception: return "(Audio not clear)"
def extract_text_from_video(f):
    clip = VideoFileClip(f); clip.audio.write_audiofile("temp.wav", verbose=False, logger=None)
    return extract_text_from_audio(open("temp.wav","rb"))
def extract_text_from_file(f):
    ext = f.name.split(".")[-1].lower()
    try:
        if ext=="pdf": pdf=fitz.open(stream=f.read(),filetype="pdf"); return "".join([p.get_text("text") for p in pdf])
        elif ext=="txt": return f.read().decode("utf-8")
        elif ext=="docx": d=docx.Document(f); return "\n".join([p.text for p in d.paragraphs])
        elif ext in ["csv","xlsx"]: df=pd.read_csv(f) if ext=="csv" else pd.read_excel(f); return df.to_string(index=False)
        elif ext=="html": return BeautifulSoup(f.read(),"html.parser").get_text()
        elif ext=="pptx": prs=Presentation(f); t=[]; [t.append(sh.text) for sl in prs.slides for sh in sl.shapes if hasattr(sh,"text")]; return "\n".join(t)
        elif ext in ["jpg","jpeg","png"]: return extract_text_from_image(f)
        elif ext in ["mp3","wav"]: return extract_text_from_audio(f)
        elif ext in ["mp4","mkv"]: return extract_text_from_video(f)
        elif ext=="json": return json.dumps(json.load(f),indent=2)
    except Exception as e: return f"(Error reading {f.name}: {e})"
    return ""

# ---------- Ollama + Embedding ----------
def get_embedding(t):
    r = requests.post(f"{OLLAMA_URL}/api/embeddings", json={"model":"nomic-embed-text","prompt":t}).json()
    return np.array(r["embedding"], dtype="float32")

def generate_answer(ctx,q,h,pn):
    hist = "\n".join([f"Student: {x}\nTeacher: {y}" for x,y in h])
    p = f"You are an AI teacher. {pn}\n\nNotes:\n{ctx}\n\n{hist}\nStudent: {q}\nTeacher:"
    r = requests.post(f"{OLLAMA_URL}/api/generate", json={"model":"llama3.2","prompt":p}, stream=True)
    a = ""; [a:=a+json.loads(l)["response"] for l in r.iter_lines() if l]; return a.strip()

def chunk_text(t,s=500): return [t[i:i+s] for i in range(0,len(t),s)]

# ---------- ChromaDB (HTTP Mode) + FAISS Fallback ----------
try:
    client = HttpClient(host=CHROMA_HOST, port=CHROMA_PORT)
    collection = client.get_or_create_collection("mentora_kb")
    st.success("✅ Connected to ChromaDB (Docker HTTP mode).")
    USE_FAISS = False
except Exception as e:
    st.warning(f"⚠️ Could not connect to ChromaDB: {e}")
    collection = None
    USE_FAISS = True

def save_to_vector_store(chunks):
    """Save embeddings to ChromaDB or FAISS"""
    try:
        if not USE_FAISS:
            # ✅ FIXED: Safe deletion for ChromaDB >=0.5
            try:
                ids = collection.get().get("ids", [])
                if ids:
                    collection.delete(ids=ids)
            except Exception as e:
                st.warning(f"⚠️ Could not clear existing ChromaDB collection: {e}")

            # Add new chunks
            for idx, c in enumerate(chunks):
                emb = get_embedding(c["text"]).tolist()
                collection.add(
                    ids=[f"chunk_{idx}"],
                    documents=[c["text"]],
                    embeddings=[emb],
                    metadatas=[{"chunk_id": idx}],
                )
            st.success(f"✅ Saved {len(chunks)} chunks to ChromaDB.")
        else:
            e = np.vstack([get_embedding(c["text"]) for c in chunks])
            index = faiss.IndexFlatL2(e.shape[1])
            index.add(e)
            np.save("vector_store.npy", e)
            faiss.write_index(index, "faiss.index")
            st.success(f"✅ Saved {len(chunks)} chunks locally (FAISS).")
    except Exception as e:
        st.warning(f"⚠️ Vector store save failed: {e}")

def search_knowledge(query, top_k=3):
    emb = get_embedding(query).reshape(1, -1)
    try:
        if not USE_FAISS:
            results = collection.query(query_embeddings=emb.tolist(), n_results=top_k)
            docs = results.get("documents", [[]])[0]
            return "\n\n".join(docs)
        else:
            if not os.path.exists("faiss.index"): return ""
            index = faiss.read_index("faiss.index")
            e = np.load("vector_store.npy")
            D, I = index.search(emb, top_k)
            return "\n\n".join([str(e[i]) for i in I[0]])
    except Exception as e:
        st.warning(f"⚠️ Search failed: {e}")
        return ""

# ---------- Check Ollama ----------
try:
    requests.get(f"{OLLAMA_URL}/api/tags", timeout=3)
    st.success("✅ Connected to Ollama.")
except:
    st.error("❌ Ollama not running. Use `ollama serve`.")
    st.stop()

# ---------- Initialize ----------
if "chat_history" not in st.session_state: st.session_state.chat_history=[]
if "class_summary" not in st.session_state: st.session_state.class_summary=None

# ---------- Sessions ----------
st.sidebar.header("📘 Sessions")
def save_chat(s,h): json.dump(h,open(os.path.join(SESSIONS_PATH,f"{s}.json"),"w",encoding="utf-8"),indent=2)
def load_chat(s): p=os.path.join(SESSIONS_PATH,f"{s}.json"); return json.load(open(p)) if os.path.exists(p) else []
def list_sessions(): return sorted([f.replace(".json","") for f in os.listdir(SESSIONS_PATH) if f.endswith(".json")])
def delete_session(s): p=os.path.join(SESSIONS_PATH,f"{s}.json"); os.remove(p) if os.path.exists(p) else None

s=list_sessions()
n=st.sidebar.text_input("➕ New Session")
if st.sidebar.button("Create"):
    if n.strip():
        sid=n.strip().replace(" ","_")
        st.session_state.current_session=sid; st.session_state.chat_history=[]; save_chat(sid,[])
        st.sidebar.success(f"✅ Created: {sid}")
    else: st.sidebar.error("Enter session name.")
if s:
    sel=st.sidebar.selectbox("📖 Load Session",s)
    c1,c2=st.sidebar.columns(2)
    if c1.button("Load"): st.session_state.current_session=sel; st.session_state.chat_history=load_chat(sel); st.sidebar.success(f"📚 Loaded: {sel}")
    if c2.button("Delete"): delete_session(sel); st.sidebar.warning(f"🗑️ Deleted: {sel}")

# ---------- Upload & Summarize ----------
st.subheader("📤 Upload Study Materials")
f=st.file_uploader("Upload PDFs, Docs, Audio, or Videos",
                   type=["pdf","txt","docx","csv","xlsx","html","pptx","json","jpg","jpeg","png","mp3","wav","mp4","mkv"],
                   accept_multiple_files=True)
if f:
    txt=""; [txt:=txt+"\n"+extract_text_from_file(x) for x in f]
    if txt.strip():
        chunks=[{"text":c} for c in chunk_text(txt)]
        save_to_vector_store(chunks)
        st.session_state.texts=chunks
        with st.spinner("🧠 Summarizing materials..."):
            try:
                pr=f"Summarize these study materials in 5 concise bullet points:\n\n{txt[:6000]}"
                r=requests.post(f"{OLLAMA_URL}/api/generate",json={"model":"llama3.2","prompt":pr},stream=True)
                sm=""; [sm:=sm+json.loads(l)["response"] for l in r.iter_lines() if l]
                st.session_state.class_summary=sm.strip(); st.success("📘 Summary generated!")
            except Exception as e: st.warning(f"⚠️ Summary failed: {e}")
    else: st.warning("No readable text extracted.")

if st.session_state.class_summary:
    st.markdown("### 📘 Class Notes Summary")
    st.info(st.session_state.class_summary)

# ---------- Voice ----------
st.markdown("### 🎙️ Ask by Voice or Text")
voice_col1, voice_col2 = st.columns([1,3])
voice_mode = voice_col1.toggle("🎤 Enable Mic Input", False)
transcribed_text=""
audio_file = voice_col2.file_uploader("Upload your question (audio file)", type=["mp3","wav","m4a","mp4"])
if audio_file:
    with st.spinner("🧠 Transcribing with Whisper..."):
        model = load_whisper_model()
        result = model.transcribe(audio_file.name)
        transcribed_text = result["text"]
        st.success(f"🗣️ You said: “{transcribed_text}”")

# ---------- Chat ----------
if "current_session" not in st.session_state: st.warning("🧑‍🎓 Start or load a session."); st.stop()
sid = st.session_state.current_session
st.subheader(f"🧑‍🏫 Mentora Session: `{sid}` — Style: {persona}")

for q,a in st.session_state.chat_history:
    with st.chat_message("user"): st.markdown(q)
    with st.chat_message("assistant"): st.markdown(a)

q = transcribed_text or st.chat_input("Ask Mentora a question…")
if q:
    with st.chat_message("user"): st.markdown(q)
    ctx = search_knowledge(q, top_k=3)
    with st.chat_message("assistant"):
        with st.spinner("🤖 Mentora is answering…"):
            a = generate_answer(ctx,q,st.session_state.chat_history,PERSONA_STYLE[persona])
            st.markdown(a)
            with tempfile.NamedTemporaryFile(delete=False,suffix=".mp3") as tts_file:
                gTTS(a).save(tts_file.name)
                st.audio(tts_file.name, format="audio/mp3")
    st.session_state.chat_history.append((q,a)); save_chat(sid,st.session_state.chat_history)

# ---------- Analytics ----------
st.markdown("### 📊 Learning Analytics")
if st.session_state.chat_history:
    q_count=len(st.session_state.chat_history)
    avg_len=np.mean([len(a[1].split()) for a in st.session_state.chat_history])
    all_text=" ".join([a[1] for a in st.session_state.chat_history])
    words=[w.lower() for w in re.findall(r"\b\w+\b", all_text) if len(w)>4]
    top=Counter(words).most_common(5)
    df=pd.DataFrame(top,columns=["Keyword","Frequency"])
    col1,col2=st.columns(2)
    col1.metric("🗣️ Questions Asked",q_count)
    col2.metric("📄 Avg Answer Length",f"{avg_len:.1f} words")
    st.write("#### 🔍 Most Discussed Topics")
    fig=px.bar(df,x="Keyword",y="Frequency",color="Keyword",title="Top 5 Keywords Discussed")
    st.plotly_chart(fig,use_container_width=True)
else:
    st.info("📈 Analytics will appear after your first Q&A session.")

# ---------- Export ----------
st.markdown("### 🧾 Export Class Notes")
exp1,exp2=st.columns(2)
def export_md():
    fn=os.path.join(EXPORTS_PATH,f"{sid}_notes.md")
    with open(fn,"w",encoding="utf-8") as f:
        f.write(f"# Mentora Class Notes — {sid}\n\n")
        for q,a in st.session_state.chat_history:
            f.write(f"### 🧑‍🎓 {q}\n{a}\n\n")
    return fn
def export_pdf():
    fn=os.path.join(EXPORTS_PATH,f"{sid}_notes.pdf")
    pdf=FPDF(); pdf.add_page(); pdf.set_font("Arial",size=12)
    pdf.cell(200,10,txt=f"Mentora Class Notes — {sid}",ln=True,align="C")
    for q,a in st.session_state.chat_history:
        pdf.multi_cell(0,10,txt=f"Q: {q}\nA: {a}\n")
    pdf.output(fn); return fn
if exp1.button("📝 Download Markdown"):
    p=export_md(); st.success(f"✅ Saved: {p}")
    with open(p,"r",encoding="utf-8") as f: st.download_button("⬇️ Download Markdown",data=f.read(),file_name=f"{sid}_notes.md")
if exp2.button("📄 Download PDF"):
    p=export_pdf(); st.success(f"✅ Saved: {p}")
    with open(p,"rb") as f: st.download_button("⬇️ Download PDF",data=f,file_name=f"{sid}_notes.pdf",mime="application/pdf")

st.markdown('<div class="footer">© 2025 Mentora — Built by Hemanth Kumar</div>', unsafe_allow_html=True)







